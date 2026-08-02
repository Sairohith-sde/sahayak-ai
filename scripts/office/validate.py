import sys
import os
from pptx import Presentation

def validate_presentation(file_path):
    print(f"Auditing PowerPoint file: {file_path}")
    if not os.path.exists(file_path):
        print(f"Error: File not found at {file_path}")
        sys.exit(1)
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"Error: Failed to open PowerPoint presentation. {e}")
        sys.exit(1)
    errors = []
    slide_count = len(prs.slides)
    print(f"Slide Count: {slide_count}/10")
    if slide_count != 10:
        errors.append(f"Expected exactly 10 slides, found {slide_count}.")
    width = prs.slide_width.inches
    height = prs.slide_height.inches
    print(f"Slide Dimensions: {width} inches x {height} inches")
    ratio = width / height
    if abs(ratio - (16/9)) > 0.05:
         errors.append(f"Expected 16:9 widescreen layout ratio, found {ratio:.2f} ({width}x{height})")
    placeholders = ["lorem", "ipsum", "xxx", "[insert", "tbd", "todo"]
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.lower()
                for ph in placeholders:
                    if ph in text:
                        errors.append(f"Slide {i+1} contains placeholder text: {ph}")
    if errors:
        print("FAILED")
        for err in errors:
            print(err)
        sys.exit(1)
    else:
        print("PASSED")
        sys.exit(0)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python validate.py <path_to_pptx>")
        sys.exit(1)
    validate_presentation(sys.argv[1])
